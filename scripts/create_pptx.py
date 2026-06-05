#!/usr/bin/env python3
"""Create editable PPTX with one figure per slide.

Fig. 1 is decomposed into individual panel images + editable text boxes
and connector arrows for maximum editability. Figures 2-7 are inserted
as single images with separate editable title/caption text boxes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

OUT_DIR = Path(__file__).resolve().parent.parent / 'pre_submission'

FIGURES = [
    {
        'file': 'fig2_sr_curves.png',
        'title': 'Figure 2',
        'caption': (
            'SR curves for a threshold detector (A/\u03b8 = 0.3). '
            'Solid: analytical SNR. Circles: Monte Carlo (15 trials). '
            'Covariate adjustment shifts the peak rightward.'
        ),
    },
    {
        'file': 'fig3_optimal_rho.png',
        'title': 'Figure 3',
        'caption': (
            'Optimal model accuracy \u03c1* and SNR improvement. '
            '(a) \u03c1* = 0 in the SR regime (\u03c3 < \u03b8); rises steeply for \u03c3 > \u03b8. '
            '(b) SNR gain at optimal \u03c1* grows exponentially with input noise.'
        ),
    },
    {
        'file': 'fig4_detection_probability.png',
        'title': 'Figure 4',
        'caption': (
            'P_D and P_FA vs. input noise for different \u03c1 (A/\u03b8 = 0.4). '
            'Higher \u03c1 suppresses both by reducing effective noise.'
        ),
    },
    {
        'file': 'fig5_roc_comparison.png',
        'title': 'Figure 5',
        'caption': (
            'ROC curves at \u03c3/\u03b8 = 1.5 (excess noise regime, A/\u03b8 = 0.4). '
            'Higher \u03c1 improves discrimination.'
        ),
    },
    {
        'file': 'fig6_mutual_information.png',
        'title': 'Figure 6',
        'caption': (
            'Mutual information I(S; E) vs. input noise. '
            'SR peak shifts rightward with covariate adjustment (\u03c1 = 0.8).'
        ),
    },
    {
        'file': 'fig7_dvs_application.png',
        'title': 'Figure 7',
        'caption': (
            'DVS astronomical observation (EBSSA, 20 recordings). '
            '(a) ROC-AUC: Fano filter 0.866 vs temporal 0.534. '
            '(b) NRR-SPR: 71.3% noise removed, 93.9% signal preserved.'
        ),
    },
]


def _add_textbox(slide, left, top, width, height, text, font_size=Pt(12),
                 bold=False, alignment=PP_ALIGN.LEFT, color=None):
    """Helper to add an editable text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.alignment = alignment
    if color:
        p.font.color.rgb = color
    return txBox


def _add_arrow(slide, start_left, start_top, end_left, end_top):
    """Add an editable connector arrow shape."""
    # Use a line connector
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top,
        end_left, end_top
    )
    connector.line.color.rgb = RGBColor(0x44, 0x44, 0x44)
    connector.line.width = Pt(1.5)
    # Add arrowhead
    connector.end_x = end_left
    connector.end_y = end_top
    return connector


def build_fig1_slide(prs):
    """Build Figure 1 as decomposed editable components.

    Three panel images placed side by side, with separate editable text boxes
    for panel labels, descriptions, and connecting arrows.
    """
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Slide title
    _add_textbox(slide, Inches(0.3), Inches(0.15), Inches(12), Inches(0.5),
                 'Figure 1: Covariate-adjusted SR framework',
                 font_size=Pt(22), bold=True)

    # Panel dimensions
    panel_w = Inches(3.8)
    panel_h = Inches(3.0)
    panel_top = Inches(1.0)
    gap = Inches(0.25)
    left_start = Inches(0.5)

    panels = [
        {'file': 'fig1a.png', 'label': '(a)', 'title': 'Threshold detector'},
        {'file': 'fig1b.png', 'label': '(b)', 'title': 'Stochastic resonance'},
        {'file': 'fig1c.png', 'label': '(c)', 'title': 'Covariate adjustment'},
    ]

    descriptions = [
        'Signal s(t) + noise n(t)\n\u2192 threshold \u03b8 \u2192 events E(t)',
        'Event-rate modulation peaks\nat intermediate noise (\u03c3 \u2248 \u03b8)',
        'Adjustment narrows noise pdf,\nshifts operating point on SR curve',
    ]

    panel_centers = []

    for i, panel in enumerate(panels):
        left = left_start + i * (panel_w + gap)
        img_path = OUT_DIR / panel['file']

        # Panel label (editable)
        _add_textbox(slide, left, panel_top - Inches(0.35), Inches(2.5), Inches(0.35),
                     f"{panel['label']} {panel['title']}",
                     font_size=Pt(14), bold=True)

        # Panel image
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), left, panel_top, panel_w, panel_h)
        else:
            # Placeholder rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, panel_top, panel_w, panel_h)
            shape.text = f"[{panel['file']}]"

        # Description text below panel (editable)
        desc_top = panel_top + panel_h + Inches(0.1)
        _add_textbox(slide, left, desc_top, panel_w, Inches(0.8),
                     descriptions[i], font_size=Pt(10))

        # Track center for arrows
        panel_centers.append(left + panel_w / 2)

    # Connecting arrows between panels (editable shapes)
    arrow_y = panel_top + panel_h / 2
    for i in range(2):
        arrow_start = left_start + (i + 1) * (panel_w + gap) - gap + Inches(0.05)
        arrow_end = left_start + (i + 1) * (panel_w + gap) - Inches(0.05)
        # Draw right-pointing arrow between panels
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            arrow_start - Inches(0.1), arrow_y - Inches(0.15),
            gap, Inches(0.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
        shape.line.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # Overall caption (editable)
    cap_top = panel_top + panel_h + Inches(1.0)
    _add_textbox(slide, Inches(0.5), cap_top, Inches(12), Inches(1.0),
                 'FIG. 1. Conceptual schematic of covariate-adjusted SR. '
                 '(a) Threshold detector: subthreshold signal + noise \u2192 events. '
                 '(b) SR: event-rate modulation peaks at \u03c3 \u2248 \u03b8. '
                 '(c) Covariate adjustment with correlation \u03c1 narrows the residual, '
                 'equivalent to leftward shift on the SR curve.',
                 font_size=Pt(11))

    return slide


def build_figure_slide(prs, fig_info):
    """Build a standard figure slide (single image + editable text)."""
    from PIL import Image

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    img_path = OUT_DIR / fig_info['file']

    if not img_path.exists():
        print(f"  WARNING: {img_path} not found, skipping")
        return slide

    # Title (editable)
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
                 fig_info['title'], font_size=Pt(24), bold=True)

    # Image — centered, scaled to fit
    with Image.open(img_path) as img:
        img_w, img_h = img.size
    aspect = img_w / img_h

    max_w = Inches(12)
    max_h = Inches(5.0)
    if aspect > (max_w / max_h):
        w = max_w
        h = int(w / aspect)
    else:
        h = max_h
        w = int(h * aspect)

    left = int((prs.slide_width - w) / 2)
    top = Inches(0.9)
    slide.shapes.add_picture(str(img_path), left, top, w, h)

    # Caption (editable)
    cap_top = top + h + Inches(0.15)
    _add_textbox(slide, Inches(0.5), cap_top, Inches(12), Inches(1.2),
                 fig_info['caption'], font_size=Pt(12))

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Fig 1: decomposed editable slide
    build_fig1_slide(prs)

    # Figs 2-7: standard slides
    for fig_info in FIGURES:
        build_figure_slide(prs, fig_info)

    out_path = OUT_DIR / 'figures.pptx'
    prs.save(str(out_path))
    print(f"PPTX saved: {out_path}")


if __name__ == '__main__':
    main()
